import asyncio
import mimetypes
import os
from typing import List, Tuple
from fastapi import HTTPException
from pptx import Presentation
import pdfplumber
from docx import Document as DocxDocument

from image_processor.utils import get_page_images_from_pdf_async

PDF_MIME_TYPES = ["application/pdf"]
TEXT_MIME_TYPES = ["text/plain"]
POWERPOINT_TYPES = [
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
]
WORD_TYPES = [
    "application/msword",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
]
SPREADSHEET_TYPES = ["text/csv", "application/csv"]
UPLOAD_ACCEPTED_DOCUMENTS = (
    PDF_MIME_TYPES + TEXT_MIME_TYPES + POWERPOINT_TYPES + WORD_TYPES
)


class DocumentsLoader:

    def __init__(self, documents: List[str]):
        self._document_paths = documents

        self._documents: List[str] = []
        self._images: List[List[str]] = []

    @property
    def documents(self):
        return self._documents

    @property
    def images(self):
        return self._images

    async def load_documents(
        self,
        temp_dir: str,
        load_text: bool = True,
        load_images: bool = False,
    ):
        documents: List[str] = []
        images: List[str] = []

        for file_path in self._document_paths:
            if not os.path.exists(file_path):
                raise HTTPException(
                    status_code=404, detail=f"File {file_path} not found"
                )

            document = ""
            imgs = []

            mime_type = mimetypes.guess_type(file_path)[0]
            if mime_type in PDF_MIME_TYPES:
                document, imgs = await self.load_pdf(
                    file_path, load_text, load_images, temp_dir
                )
            elif mime_type in TEXT_MIME_TYPES:
                document = await self.load_text(file_path)
            elif mime_type in POWERPOINT_TYPES:
                document = self.load_powerpoint(file_path)
            elif mime_type in WORD_TYPES:
                document = self.load_msword(file_path)

            documents.append(document)
            images.append(imgs)

        self._documents = documents
        self._images = images

    async def load_pdf(
        self,
        file_path: str,
        load_text: bool,
        load_images: bool,
        temp_dir: str,
    ) -> Tuple[str, List[str]]:
        image_paths = []
        document: str = ""

        if load_text:
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    document += await asyncio.to_thread(page.extract_text)

        if load_images:
            image_paths = await get_page_images_from_pdf_async(file_path, temp_dir)

        return document, image_paths

    async def load_text(self, file_path: str) -> str:
        with open(file_path, "r") as file:
            return await asyncio.to_thread(file.read)

    def load_msword(self, file_path: str) -> str:
        document = DocxDocument(file_path)
        text = "\n".join([paragraph.text for paragraph in document.paragraphs])
        return text

    def load_powerpoint(self, file_path: str) -> str:
        presentation = Presentation(file_path)

        extracted_text = ""
        for index, slide in enumerate(presentation.slides):
            extracted_text += f"# Slide {index + 1}\n"
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        extracted_text += f"{paragraph.text}\n"
                    extracted_text += "\n"
            extracted_text += "\n\n"
        return extracted_text
