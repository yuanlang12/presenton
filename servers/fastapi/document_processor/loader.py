import mimetypes
import os
from typing import List, Tuple
from fastapi import HTTPException
from langchain_community.document_loaders import TextLoader, PDFPlumberLoader
from langchain_core.documents import Document
from langchain_text_splitters import CharacterTextSplitter, MarkdownTextSplitter
from pptx import Presentation
from docx import Document as DocxDocument

from image_processor.utils import get_page_images_from_pdf_async

PDF_MIME_TYPES = ["application/pdf"]
TEXT_MIME_TYPES = ["text/plain"]
POWERPOINT_TYPES = [
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
]
WORD_TYPES = [
    "application/msword",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
]
SPREADSHEET_TYPES = ["text/csv", "application/csv"]
UPLOAD_ACCEPTED_DOCUMENTS = (
    PDF_MIME_TYPES + TEXT_MIME_TYPES + POWERPOINT_TYPES + WORD_TYPES
)


class DocumentsLoader:

    def __init__(self, documents: List[str]):
        self._document_paths = documents

        self._documents: List[Document] = []
        self._splitted_documents: List[Document] = []
        self._images: List[List[str]] = []

        self._markdown_splitter = MarkdownTextSplitter(chunk_size=500, chunk_overlap=50)
        self._text_splitter = CharacterTextSplitter(
            separator="/n", chunk_size=500, chunk_overlap=50
        )

    @property
    def documents(self):
        return self._documents

    @property
    def splitted_documents(self):
        return self._splitted_documents

    @property
    def images(self):
        return self._images

    async def load_documents(
        self,
        temp_dir: str,
        split_documents: bool = False,
        load_markdown: bool = True,
        load_images: bool = False,
    ):
        documents: List[Document] = []
        images: List[str] = []

        splitted_documents: List[Document] = []
        for file_path in self._document_paths:
            if not os.path.exists(file_path):
                raise HTTPException(
                    status_code=404, detail=f"File {file_path} not found"
                )

            docs = []
            imgs = []

            mime_type = mimetypes.guess_type(file_path)[0]
            if mime_type in PDF_MIME_TYPES:
                docs, imgs = await self.load_pdf(
                    file_path, load_markdown, load_images, temp_dir
                )
            elif mime_type in TEXT_MIME_TYPES:
                docs = self.load_text(file_path)
            elif mime_type in POWERPOINT_TYPES:
                docs = self.load_powerpoint(file_path)
            elif mime_type in WORD_TYPES:
                docs = self.load_msword(file_path)

            documents.extend(docs)
            images.append(imgs)

            if split_documents:
                splitted_documents.extend(self.split_documents(docs, mime_type))

        self._documents = documents
        self._splitted_documents = splitted_documents
        self._images = images

    def split_documents(self, documents: List[Document], mime_type):
        return self._text_splitter.split_documents(documents)

    def clip_longer_documents(self, documents: List[Document], clip_after: int = 1200):
        for document in documents:
            document.page_content = document.page_content[:clip_after]
        return documents

    async def load_pdf(
        self,
        file_path: str,
        load_markdown: bool,
        load_images: bool,
        temp_dir: str,
    ) -> Tuple[List[Document], List[str]]:
        image_paths = []
        documents: List[Document] = []

        if load_markdown:
            loader = PDFPlumberLoader(file_path)
            docs = loader.load()
            pdf_document = Document(page_content="")
            pdf_document.metadata = docs[0].metadata
            for doc in docs:
                pdf_document.page_content += doc.page_content
            documents.append(pdf_document)

        if load_images:
            image_paths = await get_page_images_from_pdf_async(file_path, temp_dir)

        return documents, image_paths

    async def decompose_pdf_to_markdown(self, document_path: str) -> str:
        raise Exception("Not Implemented")

    def load_text(self, file_path: str) -> List[Document]:
        loader = TextLoader(file_path)
        return loader.load()

    def load_msword(self, file_path: str) -> List[Document]:
        document = DocxDocument(file_path)
        text = "\n".join([paragraph.text for paragraph in document.paragraphs])
        return [Document(page_content=text)]

    def load_powerpoint(self, file_path: str) -> List[Document]:
        presentation = Presentation(file_path)

        extracted_text = ""
        for index, slide in enumerate(presentation.slides):
            extracted_text += f"# Slide {index + 1}\n"
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        extracted_text += f"{paragraph.text}\n"
                    extracted_text += "\n"
            extracted_text += "\n\n"
        return [Document(page_content=extracted_text)]
