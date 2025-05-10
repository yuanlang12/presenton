import os
from typing import List, Optional
import uuid
from lxml import etree

from pptx import Presentation
from pptx.shapes.autoshape import Shape
from pptx.slide import Slide
from pptx.chart.data import ChartData, BubbleChartData
from pptx.chart.chart import Chart
from pptx.text.text import _Paragraph, TextFrame, Font, _Run
from pptx.enum.chart import (
    XL_CHART_TYPE,
    XL_LEGEND_POSITION,
    XL_LABEL_POSITION,
)
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from lxml.etree import fromstring, tostring
from PIL import Image

from pptx.util import Pt
from graph_processor.models import (
    BarGraphDataModel,
    BubbleChartDataModel,
    GraphTypeEnum,
    LineChartDataModel,
    PieChartDataModel,
)
from pptx.dml.color import RGBColor
from ppt_generator.models.pptx_models import (
    PptxAutoShapeBoxModel,
    PptxBoxShapeEnum,
    PptxConnectorModel,
    PptxFillModel,
    PptxFontModel,
    PptxGraphBoxModel,
    PptxParagraphModel,
    PptxPictureBoxModel,
    PptxPositionModel,
    PptxPresentationModel,
    PptxShadowModel,
    PptxSlideModel,
    PptxSpacingModel,
    PptxStrokeModel,
    PptxTextBoxModel,
    PptxTextRunModel,
)
from ppt_generator.utils import (
    clip_image,
    fit_image,
    round_image_corners,
    create_circle_image,
    change_image_color,
)

BLANK_SLIDE_LAYOUT = 6


class PptxPresentationCreator:

    def __init__(self, ppt_model: PptxPresentationModel, temp_dir: str):
        self._temp_dir = temp_dir

        self._ppt_model = ppt_model
        self._slide_models = ppt_model.slides
        # self._theme = ppt_model.theme
        # self._watermark = ppt_model.watermark

        self._ppt = Presentation()
        self._ppt.slide_width = Pt(1280)
        self._ppt.slide_height = Pt(720)

        self._slide_fill = PptxFillModel(color=ppt_model.background_color)

    def create_ppt(self):
        # self.set_presentation_theme()

        for slide_model in self._slide_models:
            # Adding global shapes to slide
            if self._ppt_model.shapes:
                slide_model.shapes.append(self._ppt_model.shapes)

            self.add_and_populate_slide(slide_model)

    def set_presentation_theme(self):
        slide_master = self._ppt.slide_master
        slide_master_part = slide_master.part

        theme_part = slide_master_part.part_related_by(RT.THEME)
        theme = fromstring(theme_part.blob)

        theme_colors = self._theme.colors.theme_color_mapping
        nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

        for color_name, hex_value in theme_colors.items():
            if color_name:
                color_element = theme.xpath(
                    f"a:themeElements/a:clrScheme/a:{color_name}/a:srgbClr",
                    namespaces=nsmap,
                )[0]
                color_element.set("val", hex_value.encode("utf-8"))

        theme_part._blob = tostring(theme)

    def add_and_populate_slide(self, slide_model: PptxSlideModel):
        slide = self._ppt.slides.add_slide(self._ppt.slide_layouts[BLANK_SLIDE_LAYOUT])

        if self._slide_fill:
            self.apply_fill_to_shape(slide.background, self._slide_fill)

        for shape_model in slide_model.shapes:
            model_type = type(shape_model)

            if model_type is PptxPictureBoxModel:
                self.add_picture(slide, shape_model)

            elif model_type is PptxAutoShapeBoxModel:
                self.add_autoshape(slide, shape_model)

            elif model_type is PptxTextBoxModel:
                self.add_textbox(slide, shape_model)

            elif model_type is PptxGraphBoxModel:
                self.add_graph(slide, shape_model)

            elif model_type is PptxConnectorModel:
                self.add_connector(slide, shape_model)

        # if self._watermark:
        # Adding watermark
        # self.add_picture(slide, self.get_watermark_box_model())

    def add_connector(self, slide: Slide, connector_model: PptxConnectorModel):
        if connector_model.thickness == 0:
            return
        connector_shape = slide.shapes.add_connector(
            connector_model.type, *connector_model.position.to_pt_xyxy()
        )
        connector_shape.line.width = Pt(connector_model.thickness)
        connector_shape.line.color.rgb = RGBColor.from_string(connector_model.color)

    def add_graph(self, slide: Slide, graph_box_model: PptxGraphBoxModel):
        chart_data = None
        chart_type = None
        graph = graph_box_model.graph
        match (graph.type):
            case GraphTypeEnum.bar:
                chart_data = self.get_bar_graph(graph.data)
                chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED

            case GraphTypeEnum.scatter:
                chart_data = self.get_scatter_graph(graph.data)
                chart_type = XL_CHART_TYPE.XY_SCATTER

            case GraphTypeEnum.bubble:
                chart_data = self.get_bubble_graph(graph.data)
                chart_type = XL_CHART_TYPE.BUBBLE

            case GraphTypeEnum.line:
                chart_data = self.get_line_graph(graph.data)
                chart_type = XL_CHART_TYPE.LINE

            case GraphTypeEnum.pie:
                chart_data = self.get_pie_graph(graph.data)
                chart_type = XL_CHART_TYPE.PIE

        if chart_data:
            chart: Chart = slide.shapes.add_chart(
                chart_type, *graph_box_model.position.to_pt_list(), chart_data
            ).chart
            self.apply_graph_styles(chart, graph_box_model)

    def apply_graph_styles(self, chart, graph_box_model: PptxGraphBoxModel):
        graph = graph_box_model.graph

        if graph.type in [GraphTypeEnum.pie, GraphTypeEnum.scatter]:
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.RIGHT
        else:
            chart.has_legend = False

        if graph_box_model.legend_font:
            self.apply_font(chart.font, graph_box_model.legend_font)

        try:
            category_axis = chart.category_axis
            if graph_box_model.category_font:
                font = category_axis.tick_labels.font
                self.apply_font(font, graph_box_model.category_font)
        except:
            print("-" * 20)
            print("Could not apply category labels style")

        try:
            value_axis = chart.value_axis
            tick_labels = value_axis.tick_labels
            if graph.postfix:
                tick_labels.number_format = f'0"{graph.postfix}"'
            if graph_box_model.value_font:
                self.apply_font(tick_labels.font, graph_box_model.value_font)
        except:
            print("-" * 20)
            print("Could not apply tick labels style")

        if graph_box_model.graph.type is GraphTypeEnum.pie:
            for plot in chart.plots:
                try:
                    plot.has_data_labels = True
                    plot.data_labels.position = (
                        XL_LABEL_POSITION.OUTSIDE_END
                        if graph_box_model.graph.type is GraphTypeEnum.bar
                        else XL_LABEL_POSITION.CENTER
                    )
                    if graph.postfix:
                        plot.data_labels.number_format = f'0"{graph.postfix}"'
                    if graph_box_model.value_font:
                        self.apply_font(
                            plot.data_labels.font,
                            (
                                graph_box_model.value_font
                                if graph_box_model.graph.type is GraphTypeEnum.bar
                                else PptxFontModel(
                                    # size=self._theme.fonts.p2,
                                    size=16,
                                    bold=True,
                                    color="ffffff",
                                )
                            ),
                        )
                except:
                    print("-" * 20)
                    print("Could not apply data labels style")

    def get_bar_graph(self, graph: BarGraphDataModel):
        chart_data = ChartData()
        chart_data.categories = graph.get_categories()
        for series in graph.series:
            chart_data.add_series(series.get_name(), series.data)
        return chart_data

    def get_bubble_graph(self, graph: BubbleChartDataModel):
        chart_data = BubbleChartData()
        for each in graph.series:
            series = chart_data.add_series(each.get_name())
            for point in each.points:
                series.add_data_point(*point.to_list())
        return chart_data

    def get_line_graph(self, graph: LineChartDataModel):
        chart_data = ChartData()
        chart_data.categories = graph.get_categories()
        for series in graph.series:
            chart_data.add_series(series.get_name(), series.data)
        return chart_data

    def get_pie_graph(self, graph: PieChartDataModel):
        chart_data = ChartData()
        chart_data.categories = graph.get_categories()
        chart_data.add_series("", graph.series[0].data)
        return chart_data

    def add_picture(self, slide: Slide, picture_model: PptxPictureBoxModel):
        image_path = picture_model.picture.path
        if (
            picture_model.clip
            or picture_model.border_radius
            or picture_model.overlay
            or picture_model.object_fit
            or picture_model.shape
        ):
            try:
                image = Image.open(image_path)
            except:
                print(f"Could not open image: {image_path}")
                return

            image = image.convert("RGBA")
            # ? Applying border radius twice to support both clip and object fit
            if picture_model.border_radius:
                image = round_image_corners(image, picture_model.border_radius)
            if picture_model.object_fit:
                image = fit_image(
                    image,
                    picture_model.position.width,
                    picture_model.position.height,
                    picture_model.object_fit,
                )
            elif picture_model.clip:
                image = clip_image(
                    image,
                    picture_model.position.width,
                    picture_model.position.height,
                )
            if picture_model.border_radius:
                image = round_image_corners(image, picture_model.border_radius)
            if picture_model.shape == PptxBoxShapeEnum.CIRCLE:
                image = create_circle_image(image)
            if picture_model.overlay:
                image = change_image_color(image, picture_model.overlay)
            image_path = os.path.join(self._temp_dir, f"{str(uuid.uuid4())}.png")
            image.save(image_path)

        margined_position = self.get_margined_position(
            picture_model.position, picture_model.margin
        )

        slide.shapes.add_picture(image_path, *margined_position.to_pt_list())

    def add_autoshape(self, slide: Slide, autoshape_box_model: PptxAutoShapeBoxModel):
        position = autoshape_box_model.position
        if autoshape_box_model.margin:
            position = self.get_margined_position(position, autoshape_box_model.margin)

        autoshape = slide.shapes.add_shape(
            autoshape_box_model.type, *position.to_pt_list()
        )

        textbox = autoshape.text_frame
        textbox.word_wrap = autoshape_box_model.text_wrap

        self.apply_fill_to_shape(autoshape, autoshape_box_model.fill)
        self.apply_margin_to_text_box(textbox, autoshape_box_model.margin)
        self.apply_stroke_to_shape(autoshape, autoshape_box_model.stroke)
        self.apply_shadow_to_shape(autoshape, autoshape_box_model.shadow)
        self.apply_border_radius_to_shape(autoshape, autoshape_box_model.border_radius)

        if autoshape_box_model.paragraphs:
            self.add_paragraphs(textbox, autoshape_box_model.paragraphs)

    def add_textbox(self, slide: Slide, textbox_model: PptxTextBoxModel):
        position = textbox_model.position
        textbox_shape = slide.shapes.add_textbox(*position.to_pt_list())
        textbox_shape.width += Pt(2)

        textbox = textbox_shape.text_frame
        textbox.word_wrap = textbox_model.text_wrap

        self.apply_fill_to_shape(textbox_shape, textbox_model.fill)
        self.apply_margin_to_text_box(textbox, textbox_model.margin)
        self.add_paragraphs(textbox, textbox_model.paragraphs)

    def add_paragraphs(
        self, textbox: TextFrame, paragraph_models: List[PptxParagraphModel]
    ):
        for index, paragraph_model in enumerate(paragraph_models):
            paragraph = textbox.add_paragraph() if index > 0 else textbox.paragraphs[0]
            self.populate_paragraph(paragraph, paragraph_model)

    def populate_paragraph(
        self, paragraph: _Paragraph, paragraph_model: PptxParagraphModel
    ):
        if paragraph_model.spacing:
            self.apply_spacing_to_paragraph(paragraph, paragraph_model.spacing)

        if paragraph_model.alignment:
            paragraph.alignment = paragraph_model.alignment

        if paragraph_model.font:
            self.apply_font_to_paragraph(paragraph, paragraph_model.font)

        text_runs = []
        if paragraph_model.text:
            text_runs = self.parse_markdown_text_to_text_runs(
                paragraph_model.font, paragraph_model.text
            )
        elif paragraph_model.text_runs:
            text_runs = paragraph_model.text_runs

        for text_run_model in text_runs:
            text_run = paragraph.add_run()
            self.populate_text_run(text_run, text_run_model)

    def parse_markdown_text_to_text_runs(self, font: PptxFontModel, text: str):
        text_runs = []
        for line in text.split("\n"):
            current_pos = 0
            while current_pos < len(line):
                # Check for bold and italic (***text***)
                if (
                    line[current_pos:].startswith("***")
                    and "***" in line[current_pos + 3 :]
                ):
                    end_pos = line.find("***", current_pos + 3)
                    text_content = line[current_pos + 3 : end_pos]
                    font_json = font.model_dump()
                    font_json["bold"] = True
                    font_json["italic"] = True
                    text_runs.append(
                        PptxTextRunModel(
                            text=text_content, font=PptxFontModel(**font_json)
                        )
                    )
                    current_pos = end_pos + 3
                # Check for bold (**text**)
                elif (
                    line[current_pos:].startswith("**")
                    and "**" in line[current_pos + 2 :]
                ):
                    end_pos = line.find("**", current_pos + 2)
                    text_content = line[current_pos + 2 : end_pos]
                    font_json = font.model_dump()
                    font_json["bold"] = True
                    text_runs.append(
                        PptxTextRunModel(
                            text=text_content, font=PptxFontModel(**font_json)
                        )
                    )
                    current_pos = end_pos + 2
                # Check for italic (*text*)
                elif (
                    line[current_pos:].startswith("__")
                    and "__" in line[current_pos + 2 :]
                ):
                    end_pos = line.find("__", current_pos + 2)
                    text_content = line[current_pos + 2 : end_pos]
                    font_json = font.model_dump()
                    font_json["italic"] = True
                    text_runs.append(
                        PptxTextRunModel(
                            text=text_content, font=PptxFontModel(**font_json)
                        )
                    )
                    current_pos = end_pos + 2
                else:
                    # Find the next formatting marker or end of line
                    next_marker = float("inf")
                    for marker in ["***", "**", "__"]:
                        pos = line.find(marker, current_pos)
                        if pos != -1:
                            next_marker = min(next_marker, pos)

                    end_pos = next_marker if next_marker != float("inf") else len(line)
                    text_content = line[current_pos:end_pos]
                    if text_content:  # Only add non-empty text
                        text_runs.append(PptxTextRunModel(text=text_content, font=font))
                    current_pos = end_pos

            # Add newline if not the last line
            if line != text.split("\n")[-1]:
                text_runs.append(PptxTextRunModel(text="\n"))

        return text_runs

    def populate_text_run(self, text_run: _Run, text_run_model: PptxTextRunModel):
        text_run.text = text_run_model.text
        if text_run_model.font:
            self.apply_font(text_run.font, text_run_model.font)

    def apply_border_radius_to_shape(self, shape: Shape, border_radius: Optional[int]):
        if not border_radius:
            return
        try:
            normalized_border_radius = Pt(border_radius) / min(
                shape.width, shape.height
            )
            shape.adjustments[0] = normalized_border_radius
        except:
            print("Could not apply border radius.")

    def apply_fill_to_shape(self, shape: Shape, fill: Optional[PptxFillModel] = None):
        if not fill:
            shape.fill.background()
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(fill.color)

    def apply_stroke_to_shape(
        self, shape: Shape, stroke: Optional[PptxStrokeModel] = None
    ):
        if not stroke or stroke.thickness == 0:
            shape.line.fill.background()
        else:
            shape.line.fill.solid()
            shape.line.fill.fore_color.rgb = RGBColor.from_string(stroke.color)
            shape.line.width = Pt(stroke.thickness)

    def apply_shadow_to_shape(
        self, shape: Shape, shadow: Optional[PptxShadowModel] = None
    ):

        # Access the XML for the shape
        sp_element = shape._element
        sp_pr = sp_element.xpath("p:spPr")[0]  # Shape properties XML element

        nsmap = sp_pr.nsmap

        # # Remove existing shadow effects if present
        effect_list = sp_pr.find("a:effectLst", namespaces=nsmap)
        if effect_list:
            old_shadow = effect_list.find("a:outerShdw")
            if old_shadow:
                effect_list.remove(
                    old_shadow, namespaces=nsmap
                )  # Remove the old shadow

        if not shadow:
            return

        if not effect_list:
            effect_list = etree.SubElement(
                sp_pr, f"{{{nsmap['a']}}}effectLst", nsmap=nsmap
            )

        outer_shadow = etree.SubElement(
            effect_list,
            f"{{{nsmap['a']}}}outerShdw",
            {
                "blurRad": f"{Pt(shadow.radius)}",
                "dir": f"{shadow.angle * 1000}",
                "dist": f"{Pt(shadow.offset)}",
                "rotWithShape": "0",
            },
            nsmap=nsmap,
        )
        color_element = etree.SubElement(
            outer_shadow,
            f"{{{nsmap['a']}}}srgbClr",
            {"val": f"{shadow.color}"},
            nsmap=nsmap,
        )
        etree.SubElement(
            color_element,
            f"{{{nsmap['a']}}}alpha",
            {"val": f"{int(shadow.opacity * 100000)}"},
            nsmap=nsmap,
        )

    def get_margined_position(
        self, position: PptxPositionModel, margin: Optional[PptxSpacingModel]
    ) -> PptxPositionModel:
        if not margin:
            return position

        left = position.left + margin.left
        top = position.top + margin.top
        width = max(position.width - margin.left - margin.right, 0)
        height = max(position.height - margin.top - margin.bottom, 0)

        return PptxPositionModel(left=left, top=top, width=width, height=height)

    def apply_margin_to_text_box(
        self, text_frame: TextFrame, margin: Optional[PptxSpacingModel]
    ) -> PptxPositionModel:
        text_frame.margin_left = Pt(margin.left if margin else 0)
        text_frame.margin_right = Pt(margin.right if margin else 0)
        text_frame.margin_top = Pt(margin.top if margin else 0)
        text_frame.margin_bottom = Pt(margin.bottom if margin else 0)

    def apply_spacing_to_paragraph(
        self, paragraph: _Paragraph, spacing: PptxSpacingModel
    ):
        paragraph.space_before = Pt(spacing.top)
        paragraph.space_after = Pt(spacing.bottom)

    def apply_font_to_paragraph(self, paragraph: _Paragraph, font: PptxFontModel):
        self.apply_font(paragraph.font, font)

    def apply_font(self, font: Font, font_model: PptxFontModel):
        font.name = font_model.name
        font.color.rgb = RGBColor.from_string(font_model.color)
        font.bold = font_model.bold
        font.italic = font_model.italic
        font.size = Pt(font_model.size)

    # def get_watermark_box_model(self):
    #     watermark_asset_path = f"assets/images/{'watermark_dark.png' if self._theme == PresentationTheme.dark else 'watermark.png'}"

    #     return PptxPictureBoxModel(
    #         position=PptxPositionModel(left=1120, top=685, width=140),
    #         clip=False,
    #         picture=PptxPictureModel(
    #             is_network=False,
    #             path=watermark_asset_path,
    #         ),
    #     )

    def save(self, path: str):
        self._ppt.save(path)
